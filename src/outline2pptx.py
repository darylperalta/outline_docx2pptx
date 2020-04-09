from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from parse_docx import parseOutline
import copy
import six
import argparse
import sys
'''
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')
'''
TEMPLATE_TEXT_IDX = 0
TEMPLATE_SPEAKER_IDX = 1
TEMPLATE_TITLE_IDX = 2
TEMPLATE_MAINPOINT_IDX = 3
TEMPLATE_TEXT_2L_IDX = 4
TEMPLATE_SUBPOINT_IDX = 5

def delete_slides(presentation, index):
        xml_slides = presentation.slides._sldIdLst  
        slides = list(xml_slides)
        xml_slides.remove(slides[index])     


def duplicate_slide(pres,index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[12]
    except:
        # blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)-1]
        # print(len(pres.slide_layouts))
        blank_slide_layout = pres.slide_layouts[0]
    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(value.reltype,
                                            value._target,
                                            value.rId)

    return copied_slide

# def try_write():
#     prs = Presentation('try.pptx')
#     print('prs', prs)
#     for slide in prs.slides:
#         print('slides: ', len(prs.slides))
#         print('slide shapes shape',len(slide.shapes))
#         for shape in slide.shapes:
#             print('paragraph', len(shape.text_frame.paragraphs))
#             if not shape.has_text_frame:
#                 continue
#             for paragraph in shape.text_frame.paragraphs:
#                 print('runs ', len(paragraph.runs))
#                 for run in paragraph.runs:
#                     # text_runs.append(run.text)
#                     print(run.text)

#     print('asdf', prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text)
#     # print(prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text)
#     prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text = 'asdfkjasfd'
#     prs.save('try2.pptx')

# def try_template():
#     occasion, theme, venue, date, author, text_book, text_verses = text_bookparseOutline('sample.docx')
#     prs = Presentation('template.pptx')
#     print('prs', prs)

#     # template: 0 - verses, 1 - 
#     for slide in prs.slides:
#         print('slides: ', len(prs.slides))
#         print('slide shapes shape',len(slide.shapes))
#         for shape in slide.shapes:
#             print('paragraph', len(shape.text_frame.paragraphs))
#             if not shape.has_text_frame:
#                 continue
#             for paragraph in shape.text_frame.paragraphs:
#                 print('runs ', len(paragraph.runs))
#                 for run in paragraph.runs:
#                     # text_runs.append(run.text)
#                     print(run.text)
#     # slide = prs.slides.add_slide(6)
#     slide_temp = duplicate_slide(prs, 1)
#     print(len(slide_temp.shapes))
#     slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = 'Some Author'
#     prs.save('out.pptx')

# def print_template():
#     prs = Presentation('template.pptx')
#     print('prs', prs)

#     # template: 0 - verses, 1 - speaker, 2 - title, 3- main point, 4 - subpoints
#     index = 0
#     for slide in prs.slides:
#         print('index', index)
#         print('slides: ', len(prs.slides))
#         print('slide shapes shape',len(slide.shapes))
#         for shape in slide.shapes:
#             print('paragraph', len(shape.text_frame.paragraphs))
#             if not shape.has_text_frame:
#                 continue
#             for paragraph in shape.text_frame.paragraphs:
#                 print('runs ', len(paragraph.runs))
#                 for run in paragraph.runs:
#                     # text_runs.append(run.text)
#                     print(run.text)
#         index += 1

# def create_speaker_slide(prs, verbose, occasion, theme, venue, date, author, text_book, text_verses):
#     slide_temp = duplicate_slide(prs, TEMPLATE_SPEAKER_IDX)
#     slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = author
# def create_pptx(verbose = False):
#     occasion, theme, venue, date, author, text_book, text_verses = parseOutline('sample.docx')
#     print('deets: ', occasion, theme, venue, date, author, text_book, text_verses)
#     prs = Presentation('template.pptx')
#     print('prs', prs)

#     if verbose == True:
#         for slide in prs.slides:
#             print('slides: ', len(prs.slides))
#             print('slide shapes shape',len(slide.shapes))
#             for shape in slide.shapes:
#                 print('paragraph', len(shape.text_frame.paragraphs))
#                 if not shape.has_text_frame:
#                     continue
#                 for paragraph in shape.text_frame.paragraphs:
#                     print('runs ', len(paragraph.runs))
#                     for run in paragraph.runs:
#                         # text_runs.append(run.text)
#                         print(run.text)
#     # slide = prs.slides.add_slide(6)

#     create_speaker_slide(prs, verbose,  occasion, theme, venue, date, author, text_book, text_verses)
#     # Create Title Slide:
#     # slide_temp = duplicate_slide(prs, 1)
#     # print(len(slide_temp.shapes))
#     # slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = 
#     prs.save('out.pptx')

    # create a class
class Outline2pptx:
    '''
    Class for creating pptx
    '''
    def __init__(self, doc_fn = '14_07_2019.docx', template_pptx = 'template.pptx', out_pptx='14_07_2019.pptx', template_id = 0, verbose = True):
        title, occasion, theme, venue, date, speaker, bibleText, bibleReading, message = parseOutline(doc_fn)
        self.title = title
        self.occasion = occasion
        self.theme  = theme
        self.venue = venue
        self.date = date
        self.speaker = speaker
        self.speaker_title = 'Senior Pastor'
        self.location = 'Tatalon, Quezon City'
        # self.text_book = text_book
        # self.text_verses = text_verses
        self.bibleText = bibleText
        self.bibleReading = bibleReading
        self.verbose = verbose
        self.prs = Presentation(template_pptx)
        self.out_pptx = out_pptx
        self.message = message
        self.template_type = '2020' #new options for new template of speaker and title slides
        self.template_id = template_id
        # self.max_char = 143
        # self.max_char = 163
        if self.template_id == 0:
            # self.max_char = 143
            # self.max_char = 163
            # self.max_char2L = 84
            self.max_char = 126
            self.max_char2L = 126
        else:
            # self.max_char = 195
            self.max_char = 193
            self.max_char2L = 122
            
    def create_pptx(self):
        if self.verbose == True:
            # print('start')

            for slide in self.prs.slides:
                # print('slides: ', len(self.prs.slides))
                # print('slide shapes shape',len(slide.shapes))
                for shape in slide.shapes:
                    # print('paragraph', len(shape.text_frame.paragraphs))
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        print('runs ', len(paragraph.runs))
                        for run in paragraph.runs:
                            # text_runs.append(run.text)
                            print(run.text)
        self.create_text_slide()
        self.create_speaker_slide()
        self.create_title_slide()
        self.create_message_slide()
        delete_slides(self.prs, 0)
        delete_slides(self.prs, 0)
        delete_slides(self.prs, 0)
        delete_slides(self.prs, 0)
        delete_slides(self.prs, 0)
        delete_slides(self.prs, 0)
        
        self.prs.save(self.out_pptx)
        # print('self date', self.date)
        # print('title', self.title)

    # slide = prs.slides.add_slide(6)
    def create_speaker_slide(self):
        slide_temp = duplicate_slide(self.prs, TEMPLATE_SPEAKER_IDX)
        # print(slide_temp.placeholders)
        # for shape in slide_temp.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        # print(len(slide_temp.shapes))
        # slide_temp.shapes[0].text = 'try'
        if self.template_type == '2020':
            # print(len(slide_temp.shapes))
            # print(len(slide_temp.shapes[1].text_frame.paragraphs[3].runs))
            # print(slide_temp.shapes[1].text_frame.paragraphs[3].runs[0].text)
            # max_len = 90 
            # print(slide_temp.shapes[1].text_frame.paragraphs[3].runs[1].text)

            # print(slide_temp.shapes[1].text_frame.paragraphs[4].runs[0].text)
            if self.template_id == 0:
                max_len_par = [90, 113]
            else:
                max_len_par = [135, 159]
            slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.speaker
            slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.speaker_title
            slide_temp.shapes[1].text_frame.paragraphs[2].runs[0].text = self.venue + ((max_len_par[0]-len(self.venue+self.occasion))*' ')  + self.occasion
            slide_temp.shapes[1].text_frame.paragraphs[3].runs[0].text =  self.location + ((max_len_par[1]-len(self.location+self.date))*' ') + self.date
            if len(slide_temp.shapes[1].text_frame.paragraphs[3].runs) > 1:
                slide_temp.shapes[1].text_frame.paragraphs[3].runs[1].text =  ''

            # slide_temp.shapes[1].text_frame.paragraphs[3].runs[0].text =  'Tatalon, Quezon City' + self.date
            # slide_temp.shapes[1].text_frame.paragraphs[3].runs[1].text =  ''
            # slide_temp.shapes[1].text_frame.paragraphs[3].runs[2].text =  ''
            # slide_temp.shapes[1].text_frame.paragraphs[3].runs[3].text =  ''
            
        
        else:
            slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.speaker
            slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.venue
            slide_temp.shapes[1].text_frame.paragraphs[2].runs[0].text = self.occasion + ', '+ self.date

    def create_title_slide(self):
        if self.template_type == '2020':
            if self.template_id == 0:
                max_len_par = [70, 94]
            else:
                max_len_par = [135, 159]
            text_books = ''
            for text_id in range(len(self.bibleText)):
                if text_id != len(self.bibleText) -1:
                    text_books = text_books + self.bibleText[text_id]['book'] + '; '
                else:
                    text_books = text_books + self.bibleText[text_id]['book']
            slide_temp = duplicate_slide(self.prs, TEMPLATE_TITLE_IDX)
            slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.title
            slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = text_books

            slide_temp.shapes[1].text_frame.paragraphs[2].runs[0].text = self.speaker + ' - '+ self.speaker_title
            slide_temp.shapes[1].text_frame.paragraphs[3].runs[0].text = self.venue + ((max_len_par[0]-len(self.venue+self.occasion))*' ') + self.occasion
            if len(slide_temp.shapes[1].text_frame.paragraphs[3].runs) > 1:
                for i in range(len(slide_temp.shapes[1].text_frame.paragraphs[3].runs)-1):
                    slide_temp.shapes[1].text_frame.paragraphs[3].runs[i+1].text =  ''
            slide_temp.shapes[1].text_frame.paragraphs[4].runs[0].text =  self.location + ((max_len_par[1]-len(self.location+self.date))*' ')+ self.date
            if len(slide_temp.shapes[1].text_frame.paragraphs[4].runs) > 1:
                for i in range(len(slide_temp.shapes[1].text_frame.paragraphs[4].runs)-1):
                    slide_temp.shapes[1].text_frame.paragraphs[4].runs[i+1].text =  ''
        else:
            text_books = ''
            for text_id in range(len(self.bibleText)):
                if text_id != len(self.bibleText) -1:
                    text_books = text_books + self.bibleText[text_id]['book'] + '; '
                else:
                    text_books = text_books + self.bibleText[text_id]['book']
            slide_temp = duplicate_slide(self.prs, TEMPLATE_TITLE_IDX)
            slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.title
            # slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.bibleText[0]['book']
            slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = text_books

            slide_temp.shapes[1].text_frame.paragraphs[2].runs[0].text = self.speaker
            slide_temp.shapes[1].text_frame.paragraphs[3].runs[0].text = self.venue
            slide_temp.shapes[1].text_frame.paragraphs[4].runs[0].text = self.occasion + ', '+ self.date

    def create_text_slide(self):
        for i in range(len(self.bibleReading)):
            for j in range(len(self.bibleReading[i]['verses'])):
                if len(self.bibleReading[i]['verses'][j]) > self.max_char2L:
                    template_idx = TEMPLATE_TEXT_IDX
                else:
                    template_idx = TEMPLATE_TEXT_2L_IDX
                slide_temp = duplicate_slide(self.prs, template_idx)
                slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.bibleReading[i]['book']
                slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.bibleReading[i]['verses'][j]
                if len(self.bibleReading[i]['verses'][j]) > self.max_char:
                    # print(self.message[i]['book'])
                    verse = self.bibleReading[i]['verses'][j]
                    verse_split = verse.split()
                    # print(verse_split)
                    # verse_words.append(verse_word)
                    verses_word = ''
                    split_idx = 0
                    for verse_idx in range(len(verse_split)):
                        verses_word += verse_split[verse_idx] + ' '
                        if len(verses_word) + len(verse_split[verse_idx+1]) > self.max_char:
                            split_idx = verse_idx
                            break
                    # print('split id',split_idx)
                    # print(verse_split[split_idx])
                    text1 = ''
                    text2 = ''
                    for verse_idx in range(split_idx):
                        text1 += verse_split[verse_idx] + ' '
                    for verse_idx in range(split_idx,len(verse_split)):
                        text2 += verse_split[verse_idx] + ' '

                    text1 = text1 + '..'
                    # print('text1', text1)
                    # print(len(text1))

                    # print('text2', text2)
                    text2 = '.. ' + text2 
                    # print(len(text2))
                    slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = text1
                    if len(text2) > self.max_char2L:
                        template_idx = TEMPLATE_TEXT_IDX
                    else:
                        template_idx = TEMPLATE_TEXT_2L_IDX
                    slide_temp = duplicate_slide(self.prs, template_idx)
                    slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.bibleReading[i]['book']
                    slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = text2
                else:
                    slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.bibleReading[i]['verses'][j]

        # for i in range(len(self.text_verses)):
        #     slide_temp = duplicate_slide(self.prs, TEMPLATE_TEXT_IDX)
        #     slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.text_book
        #     slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.text_verses[i]

        for i in range(len(self.bibleText)):
            for j in range(len(self.bibleText[i]['verses'])):
                if len(self.bibleText[i]['verses'][j]) > self.max_char2L:
                    template_idx = TEMPLATE_TEXT_IDX
                else:
                    template_idx = TEMPLATE_TEXT_2L_IDX
                slide_temp = duplicate_slide(self.prs, template_idx)
                # slide_temp = duplicate_slide(self.prs, TEMPLATE_TEXT_IDX)
                slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.bibleText[i]['book']
                slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.bibleText[i]['verses'][j]
                if len(self.bibleText[i]['verses'][j]) > self.max_char:
                    # print(self.message[i]['book'])
                    verse = self.bibleText[i]['verses'][j]
                    verse_split = verse.split()
                    # print(verse_split)
                    # verse_words.append(verse_word)
                    verses_word = ''
                    split_idx = 0
                    for verse_idx in range(len(verse_split)):
                        verses_word += verse_split[verse_idx] + ' '
                        if len(verses_word) + len(verse_split[verse_idx+1]) > self.max_char:
                            split_idx = verse_idx
                            break
                    # print('split id',split_idx)
                    # print(verse_split[split_idx])
                    text1 = ''
                    text2 = ''
                    for verse_idx in range(split_idx):
                        text1 += verse_split[verse_idx] + ' '
                    for verse_idx in range(split_idx,len(verse_split)):
                        text2 += verse_split[verse_idx] + ' '

                    text1 = text1 + '..'
                    # print('text1', text1)
                    # print(len(text1))

                    # print('text2', text2)
                    text2 = '.. ' + text2 
                    # print(len(text2))
                    slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = text1
                    if len(text2) > self.max_char2L:
                        template_idx = TEMPLATE_TEXT_IDX
                    else:
                        template_idx = TEMPLATE_TEXT_2L_IDX
                    slide_temp = duplicate_slide(self.prs, template_idx)
                    # slide_temp = duplicate_slide(self.prs, TEMPLATE_TEXT_IDX)
                    slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.bibleText[i]['book']
                    slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = text2
                else:
                    slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.bibleText[i]['verses'][j]

    def create_message_slide(self):
        for i in range(len(self.message)):
            if self.message[i]['type'] == 'bible':
                for j in range(len(self.message[i]['verses'])):
                    if len(self.message[i]['verses'][j]) > self.max_char2L:
                        template_idx = TEMPLATE_TEXT_IDX
                    else:
                        template_idx = TEMPLATE_TEXT_2L_IDX
                    slide_temp = duplicate_slide(self.prs, template_idx)
                    # slide_temp = duplicate_slide(self.prs, TEMPLATE_TEXT_IDX)
                    slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.message[i]['book']
                    # print('book', self.message[i]['book'])
                    if len(self.message[i]['verses'][j]) > self.max_char:
                        # print(self.message[i]['book'])
                        verse = self.message[i]['verses'][j]
                        verse_split = verse.split()
                        # print(verse_split)
                        # verse_words.append(verse_word)
                        verses_word = ''
                        split_idx = 0
                        for verse_idx in range(len(verse_split)):
                            verses_word += verse_split[verse_idx] + ' '
                            if len(verses_word) + len(verse_split[verse_idx+1]) > self.max_char:
                                split_idx = verse_idx
                                break
                        # print('split id',split_idx)
                        # print(verse_split[split_idx])
                        text1 = ''
                        text2 = ''
                        for verse_idx in range(split_idx):
                            text1 += verse_split[verse_idx] + ' '
                        for verse_idx in range(split_idx,len(verse_split)):
                            text2 += verse_split[verse_idx] + ' '

                        text1 = text1 + '..'
                        # print('text1', text1)
                        # print(len(text1))

                        # print('text2', text2)
                        text2 = '.. ' + text2 
                        # print(len(text2))
                        slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = text1
                        if len(text2) > self.max_char2L:
                            template_idx = TEMPLATE_TEXT_IDX
                        else:
                            template_idx = TEMPLATE_TEXT_2L_IDX
                        slide_temp = duplicate_slide(self.prs, template_idx)
                        # slide_temp = duplicate_slide(self.prs, TEMPLATE_TEXT_IDX)
                        slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.message[i]['book']
                        slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = text2
                    else:
                        slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.message[i]['verses'][j]

            elif self.message[i]['type'] == 'point':
                slide_temp = duplicate_slide(self.prs, TEMPLATE_MAINPOINT_IDX)
                # print(slide_temp.shapes[1].text_frame.auto_size)
                # slide_temp.shapes[1].text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                # slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].font.underline = True
                slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.message[i]['text']
                # print(len( slide_temp.shapes[1].text_frame.paragraphs[0].runs))
                # print(self.message[i]['text'])
                slide_temp.shapes[2].text_frame.paragraphs[0].runs[0].text = self.title



def main(args=None):
    print('Python Script to convert Preaching Outline to PPTX. Default outline document is sample.docx.')
    print('To create two versions: use --two_versions')
    # Parse arguments
    if args is None:
        args = sys.argv[1:]
    # args = parse_args(args)

    parser = argparse.ArgumentParser(description='Input File.')
    #
    parser.add_argument('--input_docx', type=str, default='sample.docx',help="Path to input outline docx file. Default = sample.docx")
    parser.add_argument('--input_template', type=str, default='template_2020.pptx',help="Path to input outline template pptx file. Default = template.pptx")
    parser.add_argument('--out_pptx', type=str, default='out.pptx',help="Path to output outline pptx file. Default = out.pptx")
    parser.add_argument('--two_versions', dest='two_versions', action='store_true', help="Creates two versions of pptx", default=True)
    parser.add_argument('--verbose', dest='verbose', action='store_true', help="Verbose")
    parser.add_argument('--input_template2', type=str, default='template2_2020.pptx',help="Path to input 2nd outline template pptx file. Default = template2.pptx")
    parser.add_argument('--out_pptx2', type=str, default='out2.pptx',help="Path to 2nd output outline pptx file. Default = out2.pptx")    

    args = parser.parse_args(args)
    print('Two versions: ',args.two_versions)
    # print('verbose', args.verbose)
    print('Converting first powerpoint.')
    out2pptx = Outline2pptx(verbose=args.verbose, doc_fn = args.input_docx, template_pptx = args.input_template, out_pptx=args.out_pptx,template_id = 0)
    out2pptx.create_pptx()
    print('Powerpoint written at ', args.out_pptx)

    if args.two_versions:
        print('Converting second powerpoint.')

        out2pptx = Outline2pptx(verbose=args.verbose, doc_fn = args.input_docx, template_pptx = args.input_template2, out_pptx=args.out_pptx2, template_id =1)
        out2pptx.create_pptx()
        print('Powerpoint written at ', args.out_pptx2)

    # create_pptx(True)

if __name__ == '__main__':
    main()