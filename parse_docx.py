from pptx import Presentation
from docx import Document
from glob import glob
# from docx import getdocumenttext

# bibleText = {
# 	"text" = ''
# 	"verses" = []
# }

bibleBooks = [
    "Genesis",
    "Gen",

    "Exodus",
    "Exo",

    "Leviticus",
    "Lev ",

    "Numbers",

    "Num ",

    "Deuteronomy",
    "Deut",

    "Joshua",
    "Judges",
    "Ruth",
    "1 Samuel",
    "2 Samuel",
    "1 Sam",
    "2 Sam",
    "1 Kings",
    "2 Kings",
    "1 Chronicles",
    "2 Chronicles",
    "1 Chron ",
    "2 Chron ",
    "Ezra",
    "Nehemiah",
    "Esther",
    "Job",
    "Psalms",
    "Ps",
    "Proverbs",
    "Prov",
    "Ecclesiastes",
    "Eccl",
    "Song of Solomon",
    "Isaiah",
    "Isa",
    "Jeremiah",
    "Jer",
    "Lamentations",
    "Ezekiel",
    "Daniel",
    "Dan",    
    "Hosea",
    "Hos",
    "Joel",
    "Amos",
    "Obadiah",
    "Jonah",
    "Micah",
    "Nahum",
    "Habakkuk",
    "Zephaniah",
    "Haggai",
    "Zechariah",
    "Malachi",
    "Matthew",
    "Matt",
    "Mark",
    "Luke",
    "John",
    "Acts",
    "Romans",
    "Rom ",
    "1 Corinthians",
    "1 Cor",
    "2 Corinthians",
    "2 Cor",    
    "Galatians",
    "Gal",
    "Ephesians",
    "Eph",

    "Philippians",
    "Phil",
    "Colossians",
    "Col",

    "1 Thessalonians",
    "2 Thessalonians",
    "1 Thess",
    "2 Thess",
    "1 Timothy",
    "2 Timothy",
    "1 Tim",
    "2 Tim",
    "Titus",
    "Philemon",
    "Hebrews",
    "Heb",

    "James",
    "1 Peter",
    "2 Peter",
    "1 Pet",
    "2 Pet",
    "1 John",
    "2 John",
    "3 John",
    "Jude",
    "Revelation",
    "Rev"

]

def checkBible(text):
	for book in bibleBooks:
		if (text in book) or (book in text[0:15]):
			return True

	return False

def checkPoint(text):
	if len(text) == 0:
		return False
	elif text[0].isupper() and text[1] ==')':
		return True
	return False


def parseOutline(doc_fn = 'sample.docx', verbose = False):
	title = ''
	occasion = ''
	theme = ''
	venue = ''
	date = ''
	author = 'Bp. Reuben Abante'
	# text_verses = []
	# text_book = ''
	temp_book = ''

	temp_verses = []
	bibleText = [] # {'book': 'verses'}
	bibleReading = [] # {'book': 'verses'}
	message = []
	doc = Document(doc_fn)
	# print(getdocumenttext(document))
	# print('sdlfasf')
	# print(len(doc.paragraphs))

	for para in doc.paragraphs:
		if 'Occasion' in para.text:
			occasion = para.text.split(':')[1].strip()
			# print(occasion)
		elif 'Theme' in para.text:
			theme = para.text.split(':')[1].strip()
		elif 'Venue' in para.text:
			venue = para.text.split(':')[1].strip()
		elif 'Date' in para.text:
			date = para.text.split(':')[1].strip()
		elif 'TITLE' in para.text:
			title = para.text.split(':')[1].strip()


		# print('para ', para.text)
	if verbose == True:
		print(occasion, theme, venue, date)

	i = 0
	# for para in doc.paragraphs:
	while i <(len(doc.paragraphs)):
		if 'Text/s' in doc.paragraphs[i].text:
			i+=1
			if i >= len(doc.paragraphs):
				break
			while(len(doc.paragraphs[i].text)==0):
				i += 1
				if i >= len(doc.paragraphs):
					break
			while(1):
				# text_book = doc.paragraphs[i].text.strip()
				# i+=1	
				# while(doc.paragraphs[i].text != 'KJV'):
				# 	while(len(doc.paragraphs[i].text)==0):
				# 		i += 1
				# 	# print('text: ', doc.paragraphs[i].text)	
				# 	# print('text: ', len(doc.paragraphs[i].text))	
				# 	if doc.paragraphs[i].text == 'KJV':
				# 		break
				# 	text_verses.append(doc.paragraphs[i].text.strip())
					
				# 	# elif len(doc.paragraphs[i].text) == 0:
				# 		# print('throw')
				# 	i += 1
				temp_book = doc.paragraphs[i].text.strip()
				# print('temp book', temp_book)
				i+=1
				if i >= len(doc.paragraphs):
					break	
				while((doc.paragraphs[i].text.strip() != 'KJV') and (doc.paragraphs[i].text.strip() != 'AMP') and (doc.paragraphs[i].text.strip() != 'BBE') and (doc.paragraphs[i].text.strip() != 'ESV') and (doc.paragraphs[i].text.strip() != 'NKJV')):
					while(len(doc.paragraphs[i].text)==0):
						i += 1
						if i >= len(doc.paragraphs):
							break
					# print('text: ', doc.paragraphs[i].text)	
					# print('text: ', len(doc.paragraphs[i].text))	
					temp_verses.append(doc.paragraphs[i].text.strip())
					# if (doc.paragraphs[i].text != 'KJV') or (doc.paragraphs[i].text != 'AMP') and (doc.paragraphs[i].text != 'BBE'):
					# 	break
					i += 1
					if i >= len(doc.paragraphs):
						break
				bibleText.append({'book':temp_book, 'verses':temp_verses})
				temp_verses = []
				i += 1
				if i >= len(doc.paragraphs):
					break
				while(len(doc.paragraphs[i].text)==0):
					i += 1
				if i >= len(doc.paragraphs):
					break
				if 'Bible Reading' in doc.paragraphs[i].text:
					break

		i+=1
	# for i in range(len(doc.paragraphs)):
	i = 0
	while(i < len(doc.paragraphs)):
		# print(len(doc.paragraphs))
		if 'Bible Reading' in doc.paragraphs[i].text:
			i+=1
			if i >= len(doc.paragraphs):
				break

			while(len(doc.paragraphs[i].text)==0):
				i += 1
				if i >= len(doc.paragraphs):
					break
			while(1):
				temp_book = doc.paragraphs[i].text.strip()
				# print('temp book', temp_book)
				i+=1
				if i >= len(doc.paragraphs):
					break	
				while((doc.paragraphs[i].text.strip() != 'KJV') and (doc.paragraphs[i].text.strip() != 'AMP') and (doc.paragraphs[i].text.strip() != 'BBE') and (doc.paragraphs[i].text.strip() != 'ESV') and (doc.paragraphs[i].text.strip() != 'NKJV')):
					while(len(doc.paragraphs[i].text)==0):
						i += 1
						if i >= len(doc.paragraphs):
							break
					# print('bible reading:: ', doc.paragraphs[i].text)	
					# print('bible reading: ', len(doc.paragraphs[i].text))	
					temp_verses.append(doc.paragraphs[i].text.strip())
					# if (doc.paragraphs[i].text != 'KJV') or (doc.paragraphs[i].text != 'AMP') and (doc.paragraphs[i].text != 'BBE'):
					# 	break
					i += 1
					if i >= len(doc.paragraphs):
						break
				bibleReading.append({'book':temp_book, 'verses':temp_verses})
				temp_verses = []
				i += 1
				if i >= len(doc.paragraphs):
					break
				while(len(doc.paragraphs[i].text)==0):
					i += 1
				if i >= len(doc.paragraphs):
					break
				if 'Song/s' in doc.paragraphs[i].text:
					break
		elif 'Song/s' in doc.paragraphs[i].text:
			break
		i+=1
		if i >= len(doc.paragraphs):
			break
		# print(doc.paragraphs[i].text)
	# print(i)
	# print()
	i = 0
	temp_verses = []
	while(i < len(doc.paragraphs)):
		# print(i)
		if 'INTRODUCTION' in doc.paragraphs[i].text:

			# print('yeah')
			while(len(doc.paragraphs[i].text)==0):
				# print('i', i)
				i += 1
				if i >= len(doc.paragraphs):
					break
			if i >= len(doc.paragraphs):
				break
			# print(doc.paragraphs[i].text)
			# print(doc.paragraphs[i+1].text)
			# print(doc.paragraphs[i+2].text)
			# print(doc.paragraphs[i+3].text)
			# print(doc.paragraphs[i+4].text)

			# break
			while(i < len(doc.paragraphs)):
				if (checkBible(doc.paragraphs[i].text.strip()) and (len(doc.paragraphs[i].text.strip()) != 0)):
					# print('true bible')
					# print(doc.paragraphs[i].text)
					# print(len(doc.paragraphs[i].text))
					# print('wjat')
					if ('(KJV)' in (doc.paragraphs[i].text)):
						temp_text = doc.paragraphs[i].text.strip()
						# print('do something')
						a = temp_text.split('(KJV)')
						if len(a[1]) != 0:
							temp_book = a[0]
							temp_verses.append(a[1])
						else:
							temp_book = a[0]
							i+=1
							while (len(doc.paragraphs[i].text)!=0):
								temp_verses.append(doc.paragraphs[i].text.strip())
								temp_verses = []
								i+=1
								if i >= len(doc.paragraphs):
									break

					else:
						temp_book = doc.paragraphs[i].text.strip()
						i+=1
						while((doc.paragraphs[i].text.strip() != 'KJV') and (doc.paragraphs[i].text.strip() != 'AMP') and (doc.paragraphs[i].text.strip() != 'BBE') and (doc.paragraphs[i].text.strip() != 'ESV') and (doc.paragraphs[i].text.strip() != 'NKJV') and (doc.paragraphs[i].text.strip() != 'CJB')):
							while(len(doc.paragraphs[i].text)==0):
								i += 1
								if i >= len(doc.paragraphs):
									break
							# print('text: ', doc.paragraphs[i].text)	
							# print('text: ', len(doc.paragraphs[i].text))	
							if i >= len(doc.paragraphs):
								break
							if doc.paragraphs[i].text.strip() == 'KJV':
								break
							temp_verses.append(doc.paragraphs[i].text.strip())
							# if (doc.paragraphs[i].text != 'KJV') or (doc.paragraphs[i].text != 'AMP') and (doc.paragraphs[i].text != 'BBE'):
							# 	break
							i += 1
							if i >= len(doc.paragraphs):
								break

					message.append({'type': 'bible', 'book':temp_book, 'verses':temp_verses})
					temp_verses = []
				elif (((doc.paragraphs[i].style.name == 'List Paragraph') or (doc.paragraphs[i].style.name == 'List Numbers'))  and (len(doc.paragraphs[i].text.strip()) != 0)):
					# print('True point')
					point = doc.paragraphs[i].text.strip()
					# print('point', point, len(doc.paragraphs[i].text.strip()))
					message.append({'type': 'point', 'text': point})

				# else:
				# 	print(doc.paragraphs[i].text.strip())
				# 	print(doc.paragraphs[i].style.name)
				i+=1
		i+=1
	# print('message:::   !!! ')
	# print(message)
	# print(len(message))

	# print('bible reading', len(bibleReading))


	# print('text:')
	# print(text_book)
	# for text in text_verses:
	# 	print(text)

	return title, occasion, theme, venue, date, author, bibleText, bibleReading, message

def parseLyrics(doc_fn = 'sample_lyrics.docx', verbose = False):

	doc = Document(doc_fn)
	
	# for para in doc.paragraphs:
	lyrics = []
	
	for i in range(len(doc.paragraphs)):
		if ((len(doc.paragraphs[i].text.strip()) != 0)):
			lyrics.append(doc.paragraphs[i].text.strip())
		i+=1

	# print('lyrics')
	# print(lyrics)

	return lyrics

def parseLyricsPPTX(lyrics_dir = 'lyrics_dir\\', verbose = True):

	# doc = Document(doc_fn)
	song_fn = glob(lyrics_dir+'*.pptx')
	# song_fn2 = glob(lyrics_dir+'*.ppt')
	# song_fn = song_fn + song_fn2
	song_lyrics = []
	for i in range(len(song_fn)):
		# print(song_fn[i])
		prs = Presentation(song_fn[i])
		lyrics = []
		if verbose == True:
	        # print('start')
			print(len(prs.slides))
			for slide in prs.slides:
				# print(len(slide.shapes))
				i = 0
				for shape in slide.shapes:
					text = ''
					if not shape.has_text_frame:
						continue
					for paragraph in shape.text_frame.paragraphs:
					    for run in paragraph.runs:
					    	# print('text', run.text)
					    	text = text + ' ' + run.text.strip()
					    	text = text.strip()
					    	# print('text', text)
					if (len(text) != 0) and (text != 'end') and (text != 'END') and (text != 'end.') and (text != 'END.') and (text != 'CHORUS:') and (text != 'CHORUS.') and (text != 'CHORUS') and (text != 'Chorus:') and (text != 'Chorus.') and (text != 'Chorus'):
						lyrics.append(text)
			song_lyrics.append(lyrics)

	# song_fn.sort()
	# print(song_fn)
	# print('song lyrics')
	# print(song_lyrics)

	return song_fn, song_lyrics

def main():
	# parseLyrics()
	parseLyricsPPTX()


if __name__ == '__main__':
	main()