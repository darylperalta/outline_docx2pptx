from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from parse_docx import parseLyricsPPTX
import copy
import six
import argparse
import sys

TEMPLATE_TITLE_IDX = 0
TEMPLATE_LYRICS_IDX = 1
TEMPLATE_END_IDX = 2

def duplicate_slide(pres,index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[12]
    except:
        # print('except')
        # blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)-1]
        # print(len(pres.slide_layouts))
        blank_slide_layout = pres.slide_layouts[6]
    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(value.reltype,
                                            value._target,
                                            value.rId)

    return copied_slide

class Lyrics2pptx:
    '''
    Class for creating pptx
    '''
    def __init__(self, doc_fn = 'sample_lyrics.docx', template_pptx = 'template_lyrics_shift.pptx', out_pptx='out_lyrics_shift.pptx', verbose = False):
        # title, occasion, theme, venue, date, speaker, text_book, text_verses, bibleReading, message = parseLyrics(doc_fn)
        self.song_fn, self.song_lyrics = parseLyricsPPTX()
        # self.title = title
        self.template_pptx = template_pptx
        self.out_pptx = out_pptx
        self.prs = Presentation(self.template_pptx)
        self.verbose = verbose
        print('song_fn')
        # print(len(self.prs.slides[0].shapes[0].text_frame.paragraphs[1].runs))
        print(self.song_fn)
        print(self.song_fn[0].split('.pptx')[0])
        if self.verbose == True:
            print('start')
            print(len(self.prs.slides))
            for slide in self.prs.slides:
            	print(len(slide.shapes))
            	# print('asd')
            	# print('shapes', print(len(slide.shapes)))
            	i = 0
            	for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    # print('i',i)
                    print('num par',len(shape.text_frame.paragraphs))
                    for paragraph in shape.text_frame.paragraphs:
                        # print('runs ', len(paragraph.runs))
                        for run in paragraph.runs:
                        	print('text', run.text)
                            # text_runs.append(run.text)
                    i += 1

    # def create_pptx(self):
    def create_pptx(self):
        song_idx = 0
        while song_idx < (len(self.song_lyrics)):
            i=0
            # print('idx', TEMPLATE_TITLE_IDX)
            while i < (len(self.song_lyrics[song_idx])):
                # print('lyric', self.lyrics[i])
                if i == 0:
                    slide_temp = duplicate_slide(self.prs, TEMPLATE_TITLE_IDX)
                    slide_temp.shapes[0].text_frame.paragraphs[0].runs[0].text = self.song_lyrics[song_idx][0]
                elif i == len(self.song_lyrics[song_idx])-1:
                    slide_temp = duplicate_slide(self.prs, TEMPLATE_END_IDX)
                    slide_temp.shapes[0].text_frame.paragraphs[0].runs[0].text = self.song_lyrics[song_idx][i]

                else:
                    slide_temp = duplicate_slide(self.prs, TEMPLATE_LYRICS_IDX)

                    # slide_temp = duplicate_slide(self.prs, TEMPLATE_LYRICS_IDX)
                    # print('i', i)
                    slide_temp.shapes[0].text_frame.paragraphs[0].runs[0].text = self.song_lyrics[song_idx][i]
                i+=1
            self.prs.save('out\\'+self.song_fn[song_idx].split('.pptx')[0] + '- left aligned.pptx')
            song_idx += 1
            self.prs = Presentation(self.template_pptx)

def main(args=None):
    # Parse arguments
    # if args is None:
    #     args = sys.argv[1:]
    # # args = parse_args(args)

    parser = argparse.ArgumentParser(description='Input File.')
    #
    parser.add_argument('--input_docx', type=str, default='sample_lyrics.docx',help="Path to input lyrics docx file. Default = sample.docx")
    parser.add_argument('--input_template', type=str, default='template_lyrics_shift.pptx',help="Path to input lyrics template pptx file. Default = template.pptx")
    parser.add_argument('--out_pptx', type=str, default='out_lyrics_shift.pptx',help="Path to input lyrics template pptx file. Default = out.pptx")
    parser.add_argument('--verbose', dest='verbose', action='store_true', help="Verbose")
    args = parser.parse_args(args)
    out2pptx = Lyrics2pptx(verbose=args.verbose, doc_fn = args.input_docx, template_pptx = args.input_template, out_pptx=args.out_pptx)
    out2pptx.create_pptx()


if __name__ == '__main__':
    main()