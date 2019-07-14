from pptx import Presentation
from parse_docx import parseOutline
import copy
import six
'''
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')
'''
TEMPLATE_TEXT_IDX = 0
TEMPLATE_SPEAKER_IDX = 1
TEMPLATE_TITLE_IDX = 2
TEMPLATE_MAINPOINT_IDX = 3
TEMPLATE_SUBPOINT_IDX = 4


def duplicate_slide(pres,index):
    template = pres.slides[index]
    try:
        blank_slide_layout = pres.slide_layouts[12]
    except:
        # blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)-1]
        # print(len(pres.slide_layouts))
        blank_slide_layout = pres.slide_layouts[0]
    copied_slide = pres.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(value.reltype,
                                            value._target,
                                            value.rId)

    return copied_slide

def try_write():
    prs = Presentation('try.pptx')
    print('prs', prs)
    for slide in prs.slides:
        print('slides: ', len(prs.slides))
        print('slide shapes shape',len(slide.shapes))
        for shape in slide.shapes:
            print('paragraph', len(shape.text_frame.paragraphs))
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print('runs ', len(paragraph.runs))
                for run in paragraph.runs:
                    # text_runs.append(run.text)
                    print(run.text)

    print('asdf', prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text)
    # print(prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text)
    prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text = 'asdfkjasfd'
    prs.save('try2.pptx')

def try_template():
    occasion, theme, venue, date, author, text_book, text_verses = text_bookparseOutline('sample.docx')
    prs = Presentation('template.pptx')
    print('prs', prs)

    # template: 0 - verses, 1 - 
    for slide in prs.slides:
        print('slides: ', len(prs.slides))
        print('slide shapes shape',len(slide.shapes))
        for shape in slide.shapes:
            print('paragraph', len(shape.text_frame.paragraphs))
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print('runs ', len(paragraph.runs))
                for run in paragraph.runs:
                    # text_runs.append(run.text)
                    print(run.text)
    # slide = prs.slides.add_slide(6)
    slide_temp = duplicate_slide(prs, 1)
    print(len(slide_temp.shapes))
    slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = 'Some Author'
    prs.save('out.pptx')

def print_template():
    prs = Presentation('template.pptx')
    print('prs', prs)

    # template: 0 - verses, 1 - speaker, 2 - title, 3- main point, 4 - subpoints
    index = 0
    for slide in prs.slides:
        print('index', index)
        print('slides: ', len(prs.slides))
        print('slide shapes shape',len(slide.shapes))
        for shape in slide.shapes:
            print('paragraph', len(shape.text_frame.paragraphs))
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                print('runs ', len(paragraph.runs))
                for run in paragraph.runs:
                    # text_runs.append(run.text)
                    print(run.text)
        index += 1

# def create_speaker_slide(prs, verbose, occasion, theme, venue, date, author, text_book, text_verses):
#     slide_temp = duplicate_slide(prs, TEMPLATE_SPEAKER_IDX)
#     slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = author
# def create_pptx(verbose = False):
#     occasion, theme, venue, date, author, text_book, text_verses = parseOutline('sample.docx')
#     print('deets: ', occasion, theme, venue, date, author, text_book, text_verses)
#     prs = Presentation('template.pptx')
#     print('prs', prs)

#     if verbose == True:
#         for slide in prs.slides:
#             print('slides: ', len(prs.slides))
#             print('slide shapes shape',len(slide.shapes))
#             for shape in slide.shapes:
#                 print('paragraph', len(shape.text_frame.paragraphs))
#                 if not shape.has_text_frame:
#                     continue
#                 for paragraph in shape.text_frame.paragraphs:
#                     print('runs ', len(paragraph.runs))
#                     for run in paragraph.runs:
#                         # text_runs.append(run.text)
#                         print(run.text)
#     # slide = prs.slides.add_slide(6)

#     create_speaker_slide(prs, verbose,  occasion, theme, venue, date, author, text_book, text_verses)
#     # Create Title Slide:
#     # slide_temp = duplicate_slide(prs, 1)
#     # print(len(slide_temp.shapes))
#     # slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = 
#     prs.save('out.pptx')

    # create a class
class Outline2pptx:
    '''
    Class for creating pptx
    '''
    def __init__(self, doc_fn = '14_07_2019.docx', template_pptx = 'template.pptx', out_pptx='14_07_2019.pptx', verbose = True):
        title, occasion, theme, venue, date, speaker, text_book, text_verses, bibleReading, message = parseOutline(doc_fn)
        self.title = title
        self.occasion = occasion
        self.theme  = theme
        self.venue = venue
        self.date = date
        self.speaker = speaker
        self.text_book = text_book
        self.text_verses = text_verses
        self.bibleReading = bibleReading
        self.verbose = verbose
        self.prs = Presentation(template_pptx)
        self.out_pptx = out_pptx
        self.message = message

    def create_pptx(self):
        if self.verbose == True:
            print('start')

            for slide in self.prs.slides:
                print('slides: ', len(self.prs.slides))
                print('slide shapes shape',len(slide.shapes))
                for shape in slide.shapes:
                    print('paragraph', len(shape.text_frame.paragraphs))
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        print('runs ', len(paragraph.runs))
                        for run in paragraph.runs:
                            # text_runs.append(run.text)
                            print(run.text)
        self.create_text_slide()
        self.create_speaker_slide()
        self.create_title_slide()
        self.create_message_slide()
        self.prs.save(self.out_pptx)
        # print('self date', self.date)
        # print('title', self.title)

    # slide = prs.slides.add_slide(6)
    def create_speaker_slide(self):
        slide_temp = duplicate_slide(self.prs, TEMPLATE_SPEAKER_IDX)
        # print(slide_temp.placeholders)
        # for shape in slide_temp.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        # print(len(slide_temp.shapes))
        # slide_temp.shapes[0].text = 'try'
        slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.speaker
        slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.venue
        slide_temp.shapes[1].text_frame.paragraphs[2].runs[0].text = self.occasion + ', '+ self.date

    def create_title_slide(self):
        slide_temp = duplicate_slide(self.prs, TEMPLATE_TITLE_IDX)
        slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.title
        slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.text_book
        slide_temp.shapes[1].text_frame.paragraphs[2].runs[0].text = self.speaker
        slide_temp.shapes[1].text_frame.paragraphs[3].runs[0].text = self.venue
        slide_temp.shapes[1].text_frame.paragraphs[4].runs[0].text = self.occasion + ', '+ self.date

    def create_text_slide(self):
        for i in range(len(self.bibleReading)):
            for j in range(len(self.bibleReading[i]['verses'])):
                slide_temp = duplicate_slide(self.prs, TEMPLATE_TEXT_IDX)
                slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.bibleReading[i]['book']
                slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.bibleReading[i]['verses'][j]

        for i in range(len(self.text_verses)):
            slide_temp = duplicate_slide(self.prs, TEMPLATE_TEXT_IDX)
            slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.text_book
            slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.text_verses[i]

    def create_message_slide(self):
        for i in range(len(self.message)):
            if self.message[i]['type'] == 'bible':
                for j in range(len(self.message[i]['verses'])):
                    slide_temp = duplicate_slide(self.prs, TEMPLATE_TEXT_IDX)
                    slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.message[i]['book']
                    slide_temp.shapes[1].text_frame.paragraphs[1].runs[0].text = self.message[i]['verses'][j]
            elif self.message[i]['type'] == 'point':
                slide_temp = duplicate_slide(self.prs, TEMPLATE_MAINPOINT_IDX)
                slide_temp.shapes[1].text_frame.paragraphs[0].runs[0].text = self.message[i]['text']
                slide_temp.shapes[2].text_frame.paragraphs[0].runs[0].text = self.title



def main():
    out2pptx = Outline2pptx(verbose=True)
    out2pptx.create_pptx()
    # create_pptx(True)

if __name__ == '__main__':
    main()